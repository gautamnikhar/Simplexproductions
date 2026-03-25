#!/usr/bin/env python3
"""
Feesback Pitch Deck Generator
Personal Branding Services for Shilen Arrow — Founder, Feesback
Color Scheme: Blue (#1B3A5C) + Beige (#F5E6D0)
Font: Montserrat
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x1B, 0x3A, 0x5C)   # #1B3A5C - primary dark blue
BLUE_MED    = RGBColor(0x2C, 0x5F, 0x8A)   # #2C5F8A - medium blue
BLUE_LIGHT  = RGBColor(0x4A, 0x90, 0xC4)   # #4A90C4 - accent blue
BEIGE       = RGBColor(0xF5, 0xE6, 0xD0)   # #F5E6D0 - warm beige
BEIGE_LIGHT = RGBColor(0xFA, 0xF3, 0xEB)   # #FAF3EB - light beige
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GOLD        = RGBColor(0xC9, 0xA9, 0x62)   # #C9A962 - gold accent
GRAY        = RGBColor(0x6B, 0x6B, 0x6B)

FONT = "Montserrat"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ────────────────────────────────────────────────────────

def add_solid_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    """Add a circle (oval) shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT, line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple styled paragraphs.
    lines = [(text, font_size, color, bold), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def add_icon_circle(slide, left, top, size, icon_text, bg_color=BLUE_LIGHT, text_color=WHITE):
    """Add a circle with an icon/emoji text inside."""
    circle = add_circle(slide, left, top, size, bg_color)
    # add text on top
    txBox = slide.shapes.add_textbox(left, top, size, size)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(size / Inches(1) * 16))
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(int(size / Inches(1) * 8))
    return circle


def add_decorative_dots(slide, left, top, cols, rows, spacing, color, size=Pt(4)):
    """Add a grid of small decorative dots."""
    for r in range(rows):
        for c in range(cols):
            x = left + c * spacing
            y = top + r * spacing
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()


def add_line_shape(slide, left, top, width, color, thickness=Pt(2)):
    """Add a horizontal line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_solid_bg(slide, BLUE_DARK)

# Left beige accent panel
add_rect(slide, Inches(0), Inches(0), Inches(5.5), H, BEIGE)

# Decorative dots on beige panel
add_decorative_dots(slide, Inches(0.4), Inches(0.4), 5, 5, Inches(0.35), BLUE_LIGHT, Pt(5))

# Diagonal accent line
add_rect(slide, Inches(5.2), Inches(0), Inches(0.08), H, GOLD)

# Beige panel text
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(4), Inches(1),
            "PERSONAL BRANDING", 14, BLUE_MED, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(1),
            "PROPOSAL", 42, BLUE_DARK, bold=True)
add_line_shape(slide, Inches(0.8), Inches(3.4), Inches(1.5), GOLD, Pt(3))
add_textbox(slide, Inches(0.8), Inches(3.7), Inches(4), Inches(1),
            "for", 16, GRAY, bold=False)
add_textbox(slide, Inches(0.8), Inches(4.1), Inches(4), Inches(1.5),
            "SHILEN ARROW", 36, BLUE_DARK, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(4), Inches(1),
            "Founder — Feesback", 16, BLUE_MED, bold=False)

# Right panel content
add_textbox(slide, Inches(6.2), Inches(1.5), Inches(6), Inches(1),
            "FEESBACK", 16, BEIGE, bold=True)
add_line_shape(slide, Inches(6.2), Inches(2.1), Inches(1), GOLD, Pt(2))

# Placeholder visual — Abstract building/real estate graphic
# Building shapes
add_rect(slide, Inches(7.0), Inches(2.5), Inches(1.2), Inches(3.5), BLUE_MED)
add_rect(slide, Inches(8.5), Inches(3.2), Inches(1.5), Inches(2.8), BLUE_LIGHT)
add_rect(slide, Inches(10.3), Inches(2.8), Inches(1.0), Inches(3.2), BEIGE)
# Windows on buildings
for row in range(5):
    add_rect(slide, Inches(7.2), Inches(2.7 + row * 0.6), Inches(0.25), Inches(0.3), BEIGE_LIGHT)
    add_rect(slide, Inches(7.7), Inches(2.7 + row * 0.6), Inches(0.25), Inches(0.3), BEIGE_LIGHT)
for row in range(4):
    add_rect(slide, Inches(8.8), Inches(3.4 + row * 0.6), Inches(0.25), Inches(0.3), WHITE)
    add_rect(slide, Inches(9.3), Inches(3.4 + row * 0.6), Inches(0.25), Inches(0.3), WHITE)
for row in range(4):
    add_rect(slide, Inches(10.5), Inches(3.0 + row * 0.7), Inches(0.2), Inches(0.35), BLUE_DARK)

# Connection nodes graphic (representing matching platform)
for i, (x, y) in enumerate([(6.5, 3.5), (6.8, 4.8), (6.0, 4.2)]):
    add_circle(slide, Inches(x), Inches(y), Inches(0.3), GOLD)

# Bottom bar
add_rect(slide, Inches(0), Inches(6.8), W, Inches(0.7), GOLD)
add_textbox(slide, Inches(0.8), Inches(6.9), Inches(12), Inches(0.5),
            "Elevating Real Estate Leadership Through Strategic Personal Branding",
            14, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ABOUT US / WHO WE ARE
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BEIGE_LIGHT)

# Top accent bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), BLUE_DARK)
add_rect(slide, Inches(0), Inches(0.12), W, Inches(0.04), GOLD)

# Section header
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.5),
            "WHO WE ARE", 12, BLUE_LIGHT, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(1),
            "Your Strategic Branding Partner", 36, BLUE_DARK, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.8), Inches(2), GOLD, Pt(3))

# Main description
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(6.5), Inches(2.5),
            "We are a boutique personal branding agency specializing in building "
            "powerful, authentic brands for founders and executives in the tech "
            "and real estate space.\n\n"
            "Our mission is to transform Shilen Arrow's expertise and vision "
            "into a magnetic personal brand that attracts top agents, buyers, "
            "and industry partnerships to Feesback.",
            16, BLACK, bold=False)

# Right side — What we bring (visual cards)
card_y = Inches(2.2)
card_items = [
    ("01", "Strategic Positioning", "Defining your unique narrative in prop-tech"),
    ("02", "Content Architecture", "LinkedIn, video, thought leadership pillars"),
    ("03", "Visual Identity", "Cohesive brand aesthetics across platforms"),
    ("04", "Community Building", "Engaged audience of agents & buyers"),
]
for idx, (num, title, desc) in enumerate(card_items):
    cy = card_y + Inches(idx * 1.15)
    card = add_rounded_rect(slide, Inches(8.0), cy, Inches(4.8), Inches(1.0), WHITE)
    # Number circle
    add_circle(slide, Inches(8.2), cy + Inches(0.15), Inches(0.6), BLUE_DARK)
    add_textbox(slide, Inches(8.2), cy + Inches(0.2), Inches(0.6), Inches(0.5),
                num, 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.0), cy + Inches(0.1), Inches(3.5), Inches(0.4),
                title, 16, BLUE_DARK, bold=True)
    add_textbox(slide, Inches(9.0), cy + Inches(0.5), Inches(3.5), Inches(0.5),
                desc, 12, GRAY, bold=False)

# Decorative dots bottom-left
add_decorative_dots(slide, Inches(0.5), Inches(6.0), 8, 3, Inches(0.3), BLUE_LIGHT, Pt(4))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY PERSONAL BRANDING FOR SHILEN
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BLUE_DARK)

# Accent
add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GOLD)

# Header
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.5),
            "THE OPPORTUNITY", 12, GOLD, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(1),
            "Why Personal Branding Matters for Shilen Arrow", 34, WHITE, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.8), Inches(2.5), GOLD, Pt(3))

# Stats row
stats = [
    ("82%", "of consumers trust\na company more when\nits founder is active\non social media"),
    ("77%", "of buyers research\nthe company founder\nbefore making a\npurchase decision"),
    ("3.5×", "more engagement\nwhen content comes\nfrom a personal\nbrand vs company"),
    ("64%", "of real estate leads\ncome from referrals\ndriven by personal\nreputation"),
]

for i, (number, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    # Card background
    add_rounded_rect(slide, x, Inches(2.3), Inches(2.8), Inches(2.5), BLUE_MED)
    # Big number
    add_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(1),
                number, 44, GOLD, bold=True, alignment=PP_ALIGN.LEFT)
    # Description
    add_textbox(slide, x + Inches(0.3), Inches(3.4), Inches(2.2), Inches(1.3),
                label, 13, BEIGE_LIGHT, bold=False, alignment=PP_ALIGN.LEFT)

# Bottom insight
add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), BEIGE)
add_textbox(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.5),
            "THE INSIGHT", 12, BLUE_DARK, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.8),
            "Feesback connects real estate agents with customers — but trust starts with WHO is behind the platform. "
            "Shilen Arrow's personal brand becomes the trust signal that drives platform adoption by agents and buyers alike.",
            15, BLUE_DARK, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SCOPE OF WORK
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BEIGE_LIGHT)

# Top accent
add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), BLUE_DARK)
add_rect(slide, Inches(0), Inches(0.12), W, Inches(0.04), GOLD)

add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.5),
            "SCOPE OF WORK", 12, BLUE_LIGHT, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
            "What We'll Build Together", 36, BLUE_DARK, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.7), Inches(2), GOLD, Pt(3))

# Three columns of deliverables
columns = [
    ("BRAND FOUNDATION", [
        "Brand audit & competitor analysis",
        "Personal brand positioning statement",
        "Tone of voice & messaging guide",
        "Bio & headline optimization",
        "Visual brand direction",
    ]),
    ("CONTENT STRATEGY", [
        "Content pillar development",
        "Monthly content calendar",
        "LinkedIn optimization & strategy",
        "Thought leadership articles",
        "Video content direction",
    ]),
    ("GROWTH & VISIBILITY", [
        "Engagement strategy playbook",
        "PR & media outreach plan",
        "Speaking opportunity sourcing",
        "Network expansion strategy",
        "Monthly analytics & reporting",
    ]),
]

for i, (title, items) in enumerate(columns):
    x = Inches(0.8 + i * 4.1)
    # Column card
    add_rounded_rect(slide, x, Inches(2.2), Inches(3.8), Inches(4.8), WHITE)
    # Colored top strip
    strip_color = [BLUE_DARK, BLUE_MED, BLUE_LIGHT][i]
    add_rect(slide, x, Inches(2.2), Inches(3.8), Inches(0.55), strip_color)
    # Number badge
    add_circle(slide, x + Inches(0.15), Inches(2.28), Inches(0.4), GOLD)
    add_textbox(slide, x + Inches(0.15), Inches(2.32), Inches(0.4), Inches(0.35),
                str(i + 1), 16, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Title on strip
    add_textbox(slide, x + Inches(0.7), Inches(2.28), Inches(3), Inches(0.45),
                title, 14, WHITE, bold=True)
    # Items
    for j, item in enumerate(items):
        iy = Inches(3.0 + j * 0.7)
        # Checkmark circle
        add_circle(slide, x + Inches(0.2), iy + Inches(0.05), Inches(0.28), BEIGE)
        add_textbox(slide, x + Inches(0.2), iy + Inches(0.02), Inches(0.28), Inches(0.28),
                    "✓", 12, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.6), iy, Inches(3), Inches(0.5),
                    item, 13, BLACK, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONTENT PILLARS
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BLUE_DARK)

# Left accent
add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GOLD)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
            "CONTENT PILLARS", 12, GOLD, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
            "Shilen Arrow's Brand Narrative", 34, WHITE, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.6), Inches(2.5), GOLD, Pt(3))

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.7),
            "Five strategic content pillars to establish authority and build trust in the real estate tech space.",
            15, BEIGE, bold=False)

pillars = [
    ("VISIONARY\nFOUNDER", "The future of real estate\ntech and how Feesback\nis shaping it"),
    ("INDUSTRY\nINSIDER", "Market trends, data insights,\nand what agents & buyers\nreally need"),
    ("BUILDER'S\nJOURNEY", "Behind-the-scenes of\ngrowing Feesback — wins,\nchallenges, lessons"),
    ("COMMUNITY\nCHAMPION", "Spotlighting agents, success\nstories, and the people\nwho make real estate human"),
    ("THOUGHT\nLEADER", "Bold takes on prop-tech,\ncustomer experience, and\nthe agent economy"),
]

for i, (title, desc) in enumerate(pillars):
    x = Inches(0.5 + i * 2.5)
    # Pillar card
    add_rounded_rect(slide, x, Inches(2.6), Inches(2.3), Inches(4.2), BLUE_MED)
    # Top accent on card
    add_rect(slide, x, Inches(2.6), Inches(2.3), Inches(0.08), GOLD)
    # Icon area
    add_circle(slide, x + Inches(0.75), Inches(2.9), Inches(0.8), BEIGE)
    icons = ["▲", "◆", "■", "●", "★"]
    add_textbox(slide, x + Inches(0.75), Inches(2.95), Inches(0.8), Inches(0.7),
                icons[i], 28, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.2), Inches(3.9), Inches(1.9), Inches(1),
                title, 15, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Separator
    add_line_shape(slide, x + Inches(0.6), Inches(4.8), Inches(1.1), GOLD, Pt(2))
    # Desc
    add_textbox(slide, x + Inches(0.15), Inches(5.0), Inches(2.0), Inches(1.5),
                desc, 11, BEIGE_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TIMELINE / ROADMAP
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BEIGE_LIGHT)

add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), BLUE_DARK)
add_rect(slide, Inches(0), Inches(0.12), W, Inches(0.04), GOLD)

add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.5),
            "ROADMAP", 12, BLUE_LIGHT, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
            "90-Day Brand Launch Plan", 36, BLUE_DARK, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.7), Inches(2), GOLD, Pt(3))

# Timeline horizontal line
add_rect(slide, Inches(1.5), Inches(3.3), Inches(10.5), Inches(0.06), BLUE_DARK)

phases = [
    ("MONTH 1", "Foundation", "Brand audit\nPositioning\nVisual identity\nProfile optimization\nContent pillar setup"),
    ("MONTH 2", "Activation", "Content calendar launch\nLinkedIn strategy go-live\nFirst thought pieces\nEngagement campaign\nVideo content kickoff"),
    ("MONTH 3", "Amplification", "PR outreach begins\nSpeaking pitches sent\nCommunity engagement\nPerformance review\nStrategy refinement"),
]

for i, (month, phase, details) in enumerate(phases):
    x = Inches(1.5 + i * 3.9)
    # Timeline dot
    add_circle(slide, x + Inches(1.2), Inches(3.1), Inches(0.45), BLUE_DARK)
    add_textbox(slide, x + Inches(1.2), Inches(3.13), Inches(0.45), Inches(0.4),
                str(i + 1), 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Phase label above
    add_textbox(slide, x, Inches(2.1), Inches(3.2), Inches(0.4),
                month, 12, BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.45), Inches(3.2), Inches(0.5),
                phase, 22, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Details card below
    add_rounded_rect(slide, x + Inches(0.1), Inches(3.8), Inches(3.0), Inches(3.0), WHITE)
    add_textbox(slide, x + Inches(0.3), Inches(3.95), Inches(2.6), Inches(2.7),
                details, 13, BLACK, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PRICING
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BLUE_DARK)

add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GOLD)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
            "INVESTMENT", 12, GOLD, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
            "Choose Your Growth Plan", 36, WHITE, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.6), Inches(2.5), GOLD, Pt(3))

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.7),
            "Two tailored packages designed for maximum brand impact and ROI.",
            15, BEIGE, bold=False)

# Package 1 — Essential ($1,200)
x1 = Inches(1.5)
add_rounded_rect(slide, x1, Inches(2.5), Inches(5.0), Inches(4.5), BLUE_MED)
add_rect(slide, x1, Inches(2.5), Inches(5.0), Inches(1.2), BEIGE)
add_textbox(slide, x1 + Inches(0.4), Inches(2.6), Inches(4), Inches(0.4),
            "ESSENTIAL", 14, BLUE_DARK, bold=True)
add_textbox(slide, x1 + Inches(0.4), Inches(2.95), Inches(4), Inches(0.8),
            "$1,200", 40, BLUE_DARK, bold=True)
add_textbox(slide, x1 + Inches(2.5), Inches(3.15), Inches(2), Inches(0.4),
            "/ month", 14, GRAY, bold=False)

essential_items = [
    "Brand audit & positioning",
    "Profile optimization (LinkedIn)",
    "Content pillar strategy",
    "8 content pieces / month",
    "Monthly content calendar",
    "Basic analytics report",
    "Email support",
]
for j, item in enumerate(essential_items):
    iy = Inches(3.9 + j * 0.4)
    add_textbox(slide, x1 + Inches(0.4), iy, Inches(4.2), Inches(0.4),
                "→  " + item, 13, WHITE, bold=False)

# Package 2 — Premium ($1,500) — HIGHLIGHTED
x2 = Inches(7.0)
# Outer glow effect via slightly larger rect
add_rounded_rect(slide, x2 - Inches(0.05), Inches(2.3), Inches(5.1), Inches(4.9), GOLD)
add_rounded_rect(slide, x2, Inches(2.35), Inches(5.0), Inches(4.8), WHITE)

# "RECOMMENDED" badge
add_rounded_rect(slide, x2 + Inches(1.2), Inches(2.1), Inches(2.5), Inches(0.4), GOLD)
add_textbox(slide, x2 + Inches(1.2), Inches(2.12), Inches(2.5), Inches(0.35),
            "★  RECOMMENDED", 11, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, x2 + Inches(0.4), Inches(2.6), Inches(4), Inches(0.4),
            "PREMIUM", 14, BLUE_DARK, bold=True)
add_textbox(slide, x2 + Inches(0.4), Inches(2.95), Inches(4), Inches(0.8),
            "$1,500", 40, BLUE_DARK, bold=True)
add_textbox(slide, x2 + Inches(2.5), Inches(3.15), Inches(2), Inches(0.4),
            "/ month", 14, GRAY, bold=False)

premium_items = [
    "Everything in Essential, plus:",
    "12 content pieces / month",
    "Video content strategy & scripts",
    "Thought leadership ghostwriting",
    "PR & media outreach",
    "Speaking opportunity sourcing",
    "Weekly strategy calls",
    "Priority support",
]
for j, item in enumerate(premium_items):
    iy = Inches(3.9 + j * 0.4)
    text_color = BLUE_DARK
    add_textbox(slide, x2 + Inches(0.4), iy, Inches(4.2), Inches(0.4),
                "→  " + item, 13, text_color, bold=(j == 0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHY US / DIFFERENTIATORS
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BEIGE_LIGHT)

add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), BLUE_DARK)
add_rect(slide, Inches(0), Inches(0.12), W, Inches(0.04), GOLD)

add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.5),
            "WHY CHOOSE US", 12, BLUE_LIGHT, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
            "What Sets Us Apart", 36, BLUE_DARK, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.7), Inches(2), GOLD, Pt(3))

diff_items = [
    ("Real Estate Focus", "We understand the real estate and prop-tech landscape inside out. "
     "Your brand strategy won't be generic — it'll speak directly to agents and home buyers."),
    ("Founder-First Approach", "We don't build company brands. We build people brands. "
     "Everything is crafted around Shilen's unique voice, story, and vision."),
    ("Data-Driven Creative", "Every piece of content is backed by engagement data and audience "
     "insights. We optimize what works and pivot fast on what doesn't."),
    ("Full-Service Execution", "Strategy is just the start. We create, publish, engage, and report — "
     "so Shilen can focus on building Feesback while we build the brand."),
]

for i, (title, desc) in enumerate(diff_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(2.2 + row * 2.4)
    # Card
    add_rounded_rect(slide, x, y, Inches(5.8), Inches(2.1), WHITE)
    # Left accent bar
    add_rect(slide, x, y, Inches(0.08), Inches(2.1), GOLD)
    # Number
    add_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.55), BLUE_DARK)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.33), Inches(0.55), Inches(0.5),
                str(i + 1), 20, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(1.1), y + Inches(0.25), Inches(4.4), Inches(0.5),
                title, 18, BLUE_DARK, bold=True)
    # Description
    add_textbox(slide, x + Inches(1.1), y + Inches(0.8), Inches(4.4), Inches(1.2),
                desc, 13, GRAY, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EXPECTED OUTCOMES / KPIs
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BLUE_DARK)

add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GOLD)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
            "PROJECTED OUTCOMES", 12, GOLD, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
            "What Success Looks Like", 34, WHITE, bold=True)
add_line_shape(slide, Inches(0.8), Inches(1.6), Inches(2.5), GOLD, Pt(3))

kpis = [
    ("5,000+", "LinkedIn\nFollowers", "Targeted growth in\n90 days"),
    ("3×", "Engagement\nRate", "Above industry\naverage"),
    ("15+", "Media\nMentions", "PR & thought\nleadership hits"),
    ("50%", "Inbound\nGrowth", "More organic leads\nto Feesback"),
]

for i, (metric, label, desc) in enumerate(kpis):
    x = Inches(0.8 + i * 3.1)
    # Card
    add_rounded_rect(slide, x, Inches(2.2), Inches(2.8), Inches(3.3), BLUE_MED)
    # Metric circle
    add_circle(slide, x + Inches(0.6), Inches(2.5), Inches(1.6), BEIGE)
    add_textbox(slide, x + Inches(0.6), Inches(2.7), Inches(1.6), Inches(1),
                metric, 36, BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, x + Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.8),
                label, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.2), Inches(5.0), Inches(2.4), Inches(0.6),
                desc, 12, BEIGE_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

# Bottom bar with key message
add_rounded_rect(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), BEIGE)
add_textbox(slide, Inches(1.2), Inches(6.05), Inches(11), Inches(0.7),
            "We don't just promise results — we track, measure, and optimize every step. "
            "Monthly reports keep you informed and confident in the ROI of your brand investment.",
            14, BLUE_DARK, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — NEXT STEPS / CTA
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BEIGE_LIGHT)

# Full-width decorative header area
add_rect(slide, Inches(0), Inches(0), W, Inches(3.0), BLUE_DARK)
add_rect(slide, Inches(0), Inches(3.0), W, Inches(0.06), GOLD)

# Decorative dots
add_decorative_dots(slide, Inches(10), Inches(0.5), 6, 5, Inches(0.3), BLUE_MED, Pt(5))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
            "NEXT STEPS", 12, GOLD, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(1),
            "Let's Build Something\nExtraordinary", 40, WHITE, bold=True)

# Steps
steps = [
    ("01", "Discovery Call", "Let's dive deep into Shilen's vision, goals, and brand aspirations for Feesback."),
    ("02", "Strategy Proposal", "We'll deliver a customized brand strategy and content roadmap within 5 business days."),
    ("03", "Kick-Off", "Once approved, we hit the ground running with brand foundation work in Week 1."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 4.1)
    y = Inches(3.5)
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.0), WHITE)
    add_rect(slide, x, y, Inches(3.8), Inches(0.06), GOLD)
    add_circle(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.55), BLUE_DARK)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.33), Inches(0.55), Inches(0.5),
                num, 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.9), y + Inches(0.3), Inches(2.7), Inches(0.5),
                title, 18, BLUE_DARK, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.4), Inches(0.9),
                desc, 13, GRAY, bold=False)

# Contact / CTA area
add_rounded_rect(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(1.0), BLUE_DARK)
add_textbox(slide, Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.4),
            "Ready to elevate the Feesback brand?", 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.4),
            "Let's schedule a call  →  [Your Email]  |  [Your Phone]",
            13, BEIGE, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BLUE_DARK)

# Full decorative layout
add_rect(slide, Inches(0), Inches(0), Inches(5.5), H, BEIGE)
add_rect(slide, Inches(5.2), Inches(0), Inches(0.08), H, GOLD)

# Decorative elements on beige
add_decorative_dots(slide, Inches(0.4), Inches(5.5), 5, 4, Inches(0.35), BLUE_LIGHT, Pt(4))

# Building silhouette (mirroring slide 1)
add_rect(slide, Inches(1.0), Inches(3.5), Inches(1.0), Inches(3.0), BLUE_MED)
add_rect(slide, Inches(2.3), Inches(4.0), Inches(1.2), Inches(2.5), BLUE_LIGHT)
add_rect(slide, Inches(3.8), Inches(3.2), Inches(0.8), Inches(3.3), RGBColor(0xD4, 0xC5, 0xAD))
for row in range(4):
    add_rect(slide, Inches(1.2), Inches(3.7 + row * 0.6), Inches(0.2), Inches(0.3), BEIGE_LIGHT)
    add_rect(slide, Inches(1.6), Inches(3.7 + row * 0.6), Inches(0.2), Inches(0.3), BEIGE_LIGHT)
for row in range(3):
    add_rect(slide, Inches(2.5), Inches(4.2 + row * 0.6), Inches(0.2), Inches(0.3), WHITE)
    add_rect(slide, Inches(2.9), Inches(4.2 + row * 0.6), Inches(0.2), Inches(0.3), WHITE)

# Right side
add_textbox(slide, Inches(6.2), Inches(2.0), Inches(6), Inches(0.5),
            "THANK YOU", 48, WHITE, bold=True)
add_line_shape(slide, Inches(6.2), Inches(2.8), Inches(2), GOLD, Pt(3))

add_textbox(slide, Inches(6.2), Inches(3.2), Inches(6), Inches(1.5),
            "We're excited about the opportunity to build\n"
            "Shilen Arrow's personal brand and drive\n"
            "Feesback's growth through strategic branding.",
            16, BEIGE, bold=False)

add_textbox(slide, Inches(6.2), Inches(4.8), Inches(6), Inches(0.5),
            "Prepared exclusively for", 13, GRAY, bold=False)
add_textbox(slide, Inches(6.2), Inches(5.2), Inches(6), Inches(0.6),
            "SHILEN ARROW", 28, WHITE, bold=True)
add_textbox(slide, Inches(6.2), Inches(5.7), Inches(6), Inches(0.5),
            "Founder, Feesback", 16, GOLD, bold=False)

# Bottom gold bar
add_rect(slide, Inches(0), Inches(7.0), W, Inches(0.5), GOLD)
add_textbox(slide, Inches(0.8), Inches(7.05), Inches(12), Inches(0.4),
            "Confidential  |  Prepared with ♦ for Feesback",
            12, BLUE_DARK, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════

output_path = os.path.join(os.path.dirname(__file__), "Feesback_Personal_Branding_Pitch_Deck.pptx")
prs.save(output_path)
print(f"✅ Pitch deck saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
